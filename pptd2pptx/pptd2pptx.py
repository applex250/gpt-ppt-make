#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
pptd2pptx — 通用 PPTD → PPTX 本地渲染器 (路线 A)

输入: 任意 .pptd 项目目录 (主入口 + pages/ + media/)
输出: 本地 PPTX, 完全离线, 不依赖 kimi.com

支持 PPTD v2 子集:
  - 主题: colors / textStyles / tableStyles ($token 引用)
  - 元素: text (rich text <p>/<span>/<strong>/<em>/<u>/<s>/<sup>/<sub>),
          shape (rect/roundRect/ellipse 等常用形状), line, image (fit/crop),
          table (合并单元格/主题样式), icon (简化为色块占位)
  暂不支持: chart (后续可用 matplotlib 渲染), animations, customFonts

用法:
  python3 pptd2pptx.py /path/to/deck.pptd --output /path/to/deck.pptx [--debug]

验证:
  soffice --headless --convert-to pdf deck.pptx
  pdftoppm -jpeg -r 70 deck.pdf p && (视觉检查 p-*.jpg)
"""
import argparse
import html
import re
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ─────────────────────────── 常量 ───────────────────────────
DEFAULT_SIZE = (960, 540)          # PPTD 默认 16:9
FALLBACK_FONT = 'MiSans'
EA_FONT = 'MiSans'                 # 东亚字体 (PPTX 内嵌名称; 打开端 fallback)

ALIGN_MAP = {
    'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT,
    'justify': PP_ALIGN.JUSTIFY, 'distributed': PP_ALIGN.DISTRIBUTE,
}
ANCHOR_MAP = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}

SHAPE_MAP = {
    'rect': MSO_SHAPE.RECTANGLE, 'roundRect': MSO_SHAPE.ROUNDED_RECTANGLE,
    'ellipse': MSO_SHAPE.OVAL, 'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
    'diamond': MSO_SHAPE.DIAMOND, 'donut': MSO_SHAPE.DONUT,
    'star5': MSO_SHAPE.STAR_5_POINT, 'rightArrow': MSO_SHAPE.RIGHT_ARROW,
    'chevron': MSO_SHAPE.CHEVRON, 'homePlate': MSO_SHAPE.PENTAGON,
    'wedgeRectCallout': MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
    'bracePair': MSO_SHAPE.DOUBLE_BRACE,
    # P2 扩展 (常用形状)
    'leftArrow': MSO_SHAPE.LEFT_ARROW, 'upArrow': MSO_SHAPE.UP_ARROW,
    'downArrow': MSO_SHAPE.DOWN_ARROW, 'leftRightArrow': MSO_SHAPE.LEFT_RIGHT_ARROW,
    'upDownArrow': MSO_SHAPE.UP_DOWN_ARROW, 'pentagon': MSO_SHAPE.PENTAGON,
    'hexagon': MSO_SHAPE.HEXAGON, 'parallelogram': MSO_SHAPE.PARALLELOGRAM,
    'trapezoid': MSO_SHAPE.TRAPEZOID, 'cross': MSO_SHAPE.CROSS,
    'heart': MSO_SHAPE.HEART, 'cloud': MSO_SHAPE.CLOUD,
    'sun': MSO_SHAPE.SUN, 'moon': MSO_SHAPE.MOON, 'star4': MSO_SHAPE.STAR_4_POINT,
    'star8': MSO_SHAPE.STAR_8_POINT, 'star10': MSO_SHAPE.STAR_10_POINT,
    'flowChartProcess': MSO_SHAPE.FLOWCHART_PROCESS,
    'flowChartDecision': MSO_SHAPE.FLOWCHART_DECISION,
    'flowChartTerminator': MSO_SHAPE.FLOWCHART_TERMINATOR,
    'flowChartDocument': MSO_SHAPE.FLOWCHART_DOCUMENT,
    'flowChartData': MSO_SHAPE.FLOWCHART_DATA,
    'bentArrow': MSO_SHAPE.BENT_ARROW, 'curvedRightArrow': MSO_SHAPE.CURVED_RIGHT_ARROW,
    'ovalCallout': MSO_SHAPE.OVAL_CALLOUT, 'cloudCallout': MSO_SHAPE.CLOUD_CALLOUT,
    'doubleBrace': MSO_SHAPE.DOUBLE_BRACE, 'leftBrace': MSO_SHAPE.LEFT_BRACE,
    'rightBrace': MSO_SHAPE.RIGHT_BRACE, 'cube': MSO_SHAPE.CUBE,
    'cylinder': MSO_SHAPE.CAN, 'can': MSO_SHAPE.CAN,
}

# ─────────────────────────── 工具 ───────────────────────────
def hex_to_rgb(color: str):
    """支持 #RRGGBB / #RRGGBBAA / $token(已在调用前解析)。"""
    c = color.lstrip('#')
    if len(c) >= 6:
        return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
    return RGBColor(0, 0, 0)


class Theme:
    """解析后的主题: colors / textStyles / tableStyles, 支持 $token 解析。"""

    def __init__(self, theme_data: dict):
        self.colors = dict(theme_data.get('colors') or {})
        self.text_styles = {k: (v or {}) for k, v in (theme_data.get('textStyles') or {}).items()}
        self.table_styles = {k: (v or {}) for k, v in (theme_data.get('tableStyles') or {}).items()}
        # 内建默认样式 (对齐 PPTD 默认值)
        self.text_styles.setdefault('__default__', {})
        self._defaults = {
            'color': '#000000', 'fontSize': 18, 'fontFamily': FALLBACK_FONT,
            'bold': False, 'italic': False, 'lineHeight': 1.0, 'letterSpacing': 0,
            'marginTop': 0,
        }

    def resolve_color(self, color):
        if isinstance(color, str) and color.startswith('$'):
            return self.colors.get(color[1:], color)
        return color

    def resolve_style(self, style_key: str) -> dict:
        """取一个 textStyle (含默认链)。"""
        if not style_key:
            return {}
        key = style_key[1:] if style_key.startswith('$') else style_key
        base = dict(self.text_styles.get(key) or {})
        base.update({k: v for k, v in self._defaults.items() if k not in base})
        return base

    def merged_text_style(self, style_key: str, overrides: dict) -> dict:
        """style(theme) + content 字段 override + 默认值。"""
        style = self.resolve_style(style_key)
        merged = dict(style)
        for k, v in overrides.items():
            if v is not None:
                merged[k] = v
        for k, v in self._defaults.items():
            merged.setdefault(k, v)
        return merged

    def resolve_font(self, font_family):
        """fontFamily: string 或 {latin, ea}。"""
        if isinstance(font_family, dict):
            return font_family.get('latin', FALLBACK_FONT), font_family.get('ea', EA_FONT)
        return font_family, EA_FONT


def set_run_font(run, font_name, size_pt, color, bold=None, italic=None):
    """设置 run 的字体属性(含东亚字体 ea)。"""
    f = run.font
    if font_name:
        f.name = font_name
        rPr = run._r.get_or_add_rPr()
        for tag in ('a:ea', 'a:cs'):
            e = rPr.find(qn(tag))
            if e is None:
                e = rPr.makeelement(qn(tag), {})
                rPr.append(e)
            e.set('typeface', font_name)
    if size_pt:
        f.size = Pt(size_pt)
    if color:
        f.color.rgb = hex_to_rgb(color)
    if bold is not None:
        f.bold = bold
    if italic is not None:
        f.italic = italic


def set_shape_fill(shape, fill_data, theme: Theme):
    """Fill: solid / gradient / image。"""
    if not fill_data:
        shape.fill.background()
        return
    if isinstance(fill_data, str):
        fill_data = {'type': 'solid', 'color': fill_data}
    ftype = fill_data.get('type')
    if ftype == 'solid':
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(theme.resolve_color(fill_data['color']))
    elif ftype == 'gradient':
        stops = fill_data.get('stops') or []
        if len(stops) >= 2:
            try:
                gf = shape.fill.gradient()
                # python-pptx: gradient_stops 只读数量, 用 _fill 直接改颜色
                grad = shape.fill._fill.gradFill
                angle = fill_data.get('angle', 0)
                # 角度转换: PPTD 0=左→右, OOXML angle 是顺时针度数(单位 60000)
                from pptx.oxml.ns import qn as _qn
                lin = grad.find(_qn('a:lin'))
                if lin is not None:
                    lin.set('ang', str(int(angle * 60000)))
                for si, stop in enumerate(stops[:len(grad.gradient_stops)]):
                    gs = grad.gradient_stops[si]
                    gs.color.rgb = hex_to_rgb(theme.resolve_color(stop['color']))
            except Exception:
                # 降级: 取第一个 stop 颜色
                shape.fill.solid()
                shape.fill.fore_color.rgb = hex_to_rgb(theme.resolve_color(stops[0]['color']))
        else:
            shape.fill.background()
    elif ftype == 'image':
        shape.fill.background()   # 图像填充 python-pptx 不支持直接; 降级
    else:
        shape.fill.background()


def set_shape_border(shape, border_data, theme: Theme):
    if not border_data:
        shape.line.fill.background()
        return
    if isinstance(border_data, list):   # BorderSpec 数组 → 简化为第一项
        border_data = border_data[0] or border_data[1] or border_data[2] or border_data[3]
    if not border_data:
        shape.line.fill.background()
        return
    shape.line.color.rgb = hex_to_rgb(theme.resolve_color(border_data.get('color', '#000000')))
    w = border_data.get('width', 1)
    shape.line.width = Pt(w)
    style = border_data.get('style', 'solid')
    if style == 'dash':
        shape.line.dash_style = 2   # DASH
    elif style == 'dot':
        shape.line.dash_style = 3   # ROUND_DOT


def pptd_to_inches(bounds, page_w, page_h):
    """PPTD bounds [x, y, w, h] (px=pt) → 英寸。"""
    x, y, w, h = bounds
    return Inches(x / 72), Inches(y / 72), Inches(w / 72), Inches(h / 72)


# ─────────────────────────── 加载 ───────────────────────────
def load_pptd(manifest_path: Path):
    manifest_path = manifest_path.resolve()
    root = manifest_path.parent
    data = yaml.safe_load(manifest_path.read_text(encoding='utf-8'))
    if data.get('version') != 'v2':
        raise ValueError(f'仅支持 PPTD v2, 实际: {data.get("version")}')
    pages = []
    for entry in data.get('pages', []):
        page_path = (root / entry).resolve()
        page_data = yaml.safe_load(page_path.read_text(encoding='utf-8'))
        pages.append({'path': entry, 'data': page_data})
    return data, pages, root


# ─────────────────────────── 文本解析 ───────────────────────────
TAG_RE = re.compile(r'<(/?)(p|span|strong|em|u|s|sup|sub|br|a)([^>]*)>', re.I)
STYLE_ATTR_RE = re.compile(r'style\s*=\s*"([^"]*)"')
INLINE_STYLE_RE = re.compile(r'([a-z-]+)\s*:\s*([^;]+)')


def parse_rich_text(text):
    """把 rich text 解析成段落列表:
    [{'align':.., 'lineHeight':.., 'marginTop':.., 'runs':[
        {'text':.., 'bold':.., 'color':.., 'fontSize':.., 'fontFamily':..,
         'italic':.., 'underline':.., 'strike':.., 'sup':.., 'sub':..}]}, ...]
    """
    if text is None:
        return []
    if not isinstance(text, str):
        text = str(text)
    paragraphs = []
    cur = {'runs': [], 'align': None, 'lineHeight': None, 'marginTop': None}
    # 栈: 每项 (bold, italic, underline, strike, sup, sub, color, fontSize, fontFamily)
    stack = [dict(bold=False, italic=False, underline=False, strike=False,
                  sup=False, sub=False, color=None, fontSize=None, fontFamily=None)]

    def push_text(txt):
        if not txt:
            return
        st = stack[-1]
        cur['runs'].append({'text': txt, **st})

    # 拆分标签
    pos = 0
    for m in TAG_RE.finditer(text):
        push_text(text[pos:m.start()])
        pos = m.end()
        closing, tag, attrs = m.group(1), m.group(2).lower(), m.group(3)
        if tag == 'br':
            push_text('\n')
            continue
        if tag == 'p':
            if closing:
                paragraphs.append(cur)
                cur = {'runs': [], 'align': None, 'lineHeight': None, 'marginTop': None}
            else:
                style_text = STYLE_ATTR_RE.search(attrs)
                if style_text:
                    for k, v in INLINE_STYLE_RE.findall(style_text.group(1)):
                        k = k.strip().lower()
                        if k == 'text-align':
                            cur['align'] = v.strip()
                        elif k == 'line-height':
                            v = v.strip()
                            cur['lineHeight'] = float(v) if v.replace('.', '', 1).isdigit() else v
                        elif k == 'margin-top':
                            v = v.strip()
                            cur['marginTop'] = float(v.replace('px', '')) if 'px' in v else None
            continue
        if tag == 'span':
            if closing:
                if len(stack) > 1:
                    stack.pop()
            else:
                new = dict(stack[-1])
                style_text = STYLE_ATTR_RE.search(attrs)
                if style_text:
                    for k, v in INLINE_STYLE_RE.findall(style_text.group(1)):
                        k = k.strip().lower(); v = v.strip()
                        if k == 'color':
                            new['color'] = v
                        elif k == 'font-size':
                            new['fontSize'] = float(v.replace('px', ''))
                        elif k == 'font-family':
                            new['fontFamily'] = v.strip('"').strip("'")
                        elif k == 'background-color':
                            pass
                stack.append(new)
            continue
        # strong / em / u / s / sup / sub / a
        attr_map = {'strong': 'bold', 'em': 'italic', 'u': 'underline',
                    's': 'strike', 'sup': 'sup', 'sub': 'sub'}
        if tag in attr_map:
            key = attr_map[tag]
            if closing:
                if len(stack) > 1:
                    stack.pop()
            else:
                new = dict(stack[-1])
                new[key] = True
                stack.append(new)
            continue
        if tag == 'a':
            if closing:
                if len(stack) > 1:
                    stack.pop()
            else:
                new = dict(stack[-1])
                new['color'] = new['color'] or '#0563C1'
                new['underline'] = True
                stack.append(new)
            continue
    push_text(text[pos:])
    if cur['runs'] or any(cur[k] for k in ('align', 'lineHeight', 'marginTop')):
        paragraphs.append(cur)
    return paragraphs


def add_text_element(slide, bounds, content, theme: Theme, page_w, page_h):
    x, y, w, h = pptd_to_inches(bounds, page_w, page_h)
    style_key = content.get('style')
    overrides = {k: content.get(k) for k in
                 ('color', 'fontSize', 'fontFamily', 'bold', 'italic',
                  'lineHeight', 'lineHeightPx', 'letterSpacing', 'marginTop')
                 if content.get(k) is not None}
    style = theme.merged_text_style(style_key, overrides)

    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = content.get('wrap', True)

    align_v = content.get('align') or ['left', 'top']
    h_align, v_align = align_v[0], align_v[1] if len(align_v) > 1 else 'top'
    tf.vertical_anchor = ANCHOR_MAP.get(v_align, MSO_ANCHOR.TOP)
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0

    raw_text = content.get('text', '')
    if not raw_text:
        return tb
    paragraphs = parse_rich_text(raw_text)
    if not paragraphs:
        return tb

    # 若只有纯文本(无标签), 直接单段
    base_font, base_ea = theme.resolve_font(style.get('fontFamily'))
    line_h = style.get('lineHeight') or 1.0
    if isinstance(line_h, str) and 'px' in line_h:
        line_h = float(line_h.replace('px', '')) / (style.get('fontSize') or 18)

    for pi, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = ALIGN_MAP.get(para.get('align') or h_align, PP_ALIGN.LEFT)
        if para.get('lineHeight'):
            lh = para['lineHeight']
            if isinstance(lh, str) and 'px' in lh:
                p.line_spacing = Pt(float(lh.replace('px', '')))
            else:
                p.line_spacing = float(lh)
        else:
            p.line_spacing = float(line_h)
        if para.get('marginTop'):
            p.space_before = Pt(para['marginTop'])
        for run in para['runs']:
            r = p.add_run()
            r.text = run['text']
            fs = run.get('fontSize') or style.get('fontSize') or 18
            set_run_font(
                r,
                font_name=run.get('fontFamily') or base_font,
                size_pt=fs,
                color=theme.resolve_color(run.get('color') or style.get('color')),
                bold=run.get('bold', style.get('bold', False)),
                italic=run.get('italic', style.get('italic', False)),
            )
            if run.get('underline'):
                r.font.underline = True
            if run.get('strike'):
                r.font.strike = True
            if run.get('sup'):
                r.font._rPr.set('baseline', '30000')
            if run.get('sub'):
                r.font._rPr.set('baseline', '-25000')
    return tb


# ─────────────────────────── 元素渲染分发 ───────────────────────────
def render_element(slide, el, theme: Theme, page_w, page_h, root: Path, debug=False):
    etype = el.get('elementType')
    x, y, w, h = pptd_to_inches(el['bounds'], page_w, page_h)
    if etype == 'text':
        add_text_element(slide, el['bounds'], el.get('content') or {}, theme, page_w, page_h)
    elif etype == 'shape':
        _render_shape(slide, el, theme, x, y, w, h)
    elif etype == 'line':
        _render_line(slide, el, theme, x, y, w, h)
    elif etype == 'image':
        _render_image(slide, el, x, y, w, h, root)
    elif etype == 'table':
        _render_table(slide, el, theme, x, y, w, h)
    elif etype == 'icon':
        _render_icon(slide, el, theme, x, y, w, h, debug=debug)
    else:
        if debug:
            print(f'  [unknown] elementType={etype}, 跳过')


def _render_icon(slide, el, theme, x, y, w, h, debug=False):
    """Font Awesome 图标: 本地 SVG → cairosvg 渲染 → 嵌入 PPTX。
    找不到图标时降级为色块占位。
    """
    icon_name = el.get('iconName', '')
    style, _, name = icon_name.partition(':')
    if not name:
        name = style; style = 'fas'
    fill = el.get('fill') or {}
    if isinstance(fill, str):
        fill = {'type': 'solid', 'color': fill}
    color = theme.resolve_color(fill.get('color', '#000000'))
    icon_dir = Path(__file__).parent / 'assets' / 'fa-icons'
    svg_path = icon_dir / f'{name}.svg'
    if svg_path.is_file():
        try:
            import cairosvg
            svg = svg_path.read_text(encoding='utf-8')
            # 注入 fill 色到 <path> 标签 (兼容旧版 cairosvg, 不用 css 参数)
            svg = svg.replace('<path ', f'<path fill="{color}" ')
            png_bytes = cairosvg.svg2png(bytestring=svg.encode('utf-8'))
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tf:
                tf.write(png_bytes)
                tmp_path = tf.name
            slide.shapes.add_picture(tmp_path, x, y, width=w, height=h)
            import os
            os.unlink(tmp_path)
            return
        except Exception as e:
            if debug:
                print(f'  [icon] {icon_name} 渲染失败({e}), 降级色块')
    # 降级: 色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(color)
    shape.line.fill.background()
    if debug:
        print(f'  [icon] {icon_name} → 色块占位(库中无此图标)')


def _render_shape(slide, el, theme, x, y, w, h):
    name = el.get('shapeName', 'rect')
    mso = SHAPE_MAP.get(name)
    if mso is None:
        if el.get('elementId') == 'bg' or name == 'rect':
            mso = MSO_SHAPE.RECTANGLE
        else:
            mso = MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(mso, x, y, w, h)
    if name == 'roundRect':
        try:
            shape.adjustments[0] = (el.get('adjustments') or [16667])[0] / 100000
        except Exception:
            pass
    set_shape_fill(shape, el.get('fill'), theme)
    set_shape_border(shape, el.get('border'), theme)
    shape.shadow.inherit = False


def _render_line(slide, el, theme, x, y, w, h):
    points = (el.get('points') or '').strip().split()
    if len(points) < 2:
        return
    try:
        vb_w, vb_h = el.get('viewBox', [1, 1])
    except TypeError:
        vb_w, vb_h = 1, 1
    p0 = [float(v) for v in points[0].split(',')]
    p1 = [float(v) for v in points[-1].split(',')]
    # 归一化坐标 → 英寸
    x0 = x + Inches((p0[0] / vb_w) * (w / Inches(1)))
    y0 = y + Inches((p0[1] / vb_h) * (h / Inches(1)))
    x1 = x + Inches((p1[0] / vb_w) * (w / Inches(1)))
    y1 = y + Inches((p1[1] / vb_h) * (h / Inches(1)))
    conn = slide.shapes.add_connector(1, x0, y0, x1, y1)   # straight
    border = el.get('border') or {}
    color = theme.resolve_color(border.get('color', '#000000'))
    conn.line.color.rgb = hex_to_rgb(color)
    conn.line.width = Pt(border.get('width', 1))
    conn.shadow.inherit = False


def _render_image(slide, el, x, y, w, h, root):
    src = el.get('src', '')
    if src.startswith(('http://', 'https://', 'data:')):
        return   # 远程图片: 本地渲染不支持, 跳过 (保持占位)
    img_path = (root / src).resolve()
    if not img_path.is_file():
        return
    fit = (el.get('fit') or {}).get('mode', 'cover')
    try:
        from PIL import Image
        iw, ih = Image.open(img_path).size
    except Exception:
        iw, ih = 400, 300
    if fit == 'fill':
        slide.shapes.add_picture(str(img_path), x, y, width=w, height=h)
    elif fit == 'contain':
        ratio = min(w / Inches(iw / 72), h / Inches(ih / 72))
        pw = Inches(iw / 72) * ratio
        ph = Inches(ih / 72) * ratio
        slide.shapes.add_picture(str(img_path), x + (w - pw) / 2, y + (h - ph) / 2,
                                 width=pw, height=ph)
    else:   # cover
        ratio = max(w / Inches(iw / 72), h / Inches(ih / 72))
        pw = Inches(iw / 72) * ratio
        ph = Inches(ih / 72) * ratio
        pic = slide.shapes.add_picture(str(img_path), x + (w - pw) / 2, y + (h - ph) / 2,
                                       width=pw, height=ph)


def _set_cell_border(cell, border_data, theme: Theme):
    """给表格单元格设置边框 (python-pptx _Cell 无 .line, 直接写 XML)。"""
    from pptx.oxml.ns import qn as _qn
    from lxml import etree
    tcPr = cell._tc.get_or_add_tcPr()
    # 清除已有边框
    for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for e in tcPr.findall(_qn(tag)):
            tcPr.remove(e)
    if isinstance(border_data, list):   # BorderSpec 数组 [top,right,bottom,left] 或 [tb,lr]
        if len(border_data) == 4:
            sides = [('a:lnT', border_data[0]), ('a:lnR', border_data[1]),
                     ('a:lnB', border_data[2]), ('a:lnL', border_data[3])]
        else:
            sides = [('a:lnT', border_data[0]), ('a:lnB', border_data[0]),
                     ('a:lnL', border_data[1]), ('a:lnR', border_data[1])]
    else:
        sides = [('a:lnT', border_data), ('a:lnR', border_data),
                 ('a:lnB', border_data), ('a:lnL', border_data)]
    for tag, bd in sides:
        if bd is None:
            continue
        color = theme.resolve_color(bd.get('color', '#000000')) if isinstance(bd, dict) else '#000000'
        width = bd.get('width', 1) if isinstance(bd, dict) else 1
        ln = tcPr.makeelement(_qn(tag), {'w': str(int(width * 12700)), 'cap': 'flat'})
        fill = ln.makeelement(_qn('a:solidFill'), {})
        clr = fill.makeelement(_qn('a:srgbClr'), {'val': color.lstrip('#')[:6]})
        fill.append(clr)
        ln.append(fill)
        tcPr.append(ln)


def _render_table(slide, el, theme, x, y, w, h):
    rows_data = el.get('rows') or []
    if not rows_data:
        return
    n_rows = len(rows_data)
    n_cols = max(len(r) for r in rows_data)
    col_widths = el.get('columnWidths') or [1.0 / n_cols] * n_cols
    row_heights = el.get('rowHeights') or [1.0 / n_rows] * n_rows

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = int(w * cw)
    for ri, rh in enumerate(row_heights):
        tbl.rows[ri].height = int(h * rh)

    style_key = el.get('style')
    tstyle = {}
    if style_key and isinstance(style_key, str) and style_key.startswith('$'):
        tstyle = theme.table_styles.get(style_key[1:]) or {}
    elif isinstance(style_key, dict):
        tstyle = style_key
    cell_style = tstyle.get('cellStyle') or {}
    first_row_style = tstyle.get('firstRowStyle') or {}
    body_styles = tstyle.get('bodyStyles') or []

    def cell_merged_text(c):
        t = c.get('text') or ''
        return t

    grid = [[None] * n_cols for _ in range(n_rows)]
    # 处理合并: 先扫描 rowSpan/colSpan
    for ri, row in enumerate(rows_data):
        col = 0
        for cell in row:
            while col < n_cols and grid[ri][col] is not None:
                col += 1
            if col >= n_cols:
                break
            rs = cell.get('rowSpan', 1)
            cs = cell.get('colSpan', 1)
            grid[ri][col] = ('span', rs, cs)
            if rs > 1 or cs > 1:
                for i in range(ri, min(ri + rs, n_rows)):
                    for j in range(col, min(col + cs, n_cols)):
                        if (i, j) != (ri, col):
                            grid[i][j] = ('occ', 0, 0)
            col += cs

    for ri in range(n_rows):
        for ci in range(n_cols):
            g = grid[ri][ci]
            if g is None:
                continue
            kind = g[0]
            if kind == 'occ':
                continue
            cell = tbl.cell(ri, ci)
            if kind == 'span':
                rs, cs = g[1], g[2]
                if rs > 1 or cs > 1:
                    cell.merge(tbl.cell(min(ri + rs - 1, n_rows - 1),
                                        min(ci + cs - 1, n_cols - 1)))
            # 找源数据
            src_cell = None
            row_src = rows_data[ri]
            idx = 0
            for c in row_src:
                if idx == ci:
                    src_cell = c
                    break
                idx += c.get('colSpan', 1)
            if src_cell is None:
                continue
            # 样式: cell 内联 > firstRow/body 分类 > cellStyle > 默认
            st = dict(cell_style)
            if ri == 0 and first_row_style:
                st.update(first_row_style)
            if body_styles and ri > 0:
                st.update(body_styles[(ri - 1) % len(body_styles)])
            st.update({k: v for k, v in src_cell.items()
                       if k in ('color', 'fontSize', 'fontFamily', 'bold', 'italic',
                                'fill', 'border', 'align', 'lineHeight', 'backgroundColor') and v is not None})
            txt = src_cell.get('text') or ''
            cell.text = ''
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = cell.text_frame.paragraphs[0]
            align = st.get('align') or ['center', 'middle']
            para.alignment = ALIGN_MAP.get(align[0], PP_ALIGN.CENTER)
            if txt:
                runs = parse_rich_text(txt)
                if runs:
                    tf = cell.text_frame
                    tf.word_wrap = True
                    fname, fea = theme.resolve_font(st.get('fontFamily') or FALLBACK_FONT)
                    for run in runs[0]['runs']:
                        r = para.add_run()
                        r.text = run['text']
                        fs = run.get('fontSize') or st.get('fontSize') or 13
                        color = st.get('color') or '#000000'
                        set_run_font(r, font_name=run.get('fontFamily') or fname,
                                     size_pt=fs, color=theme.resolve_color(color),
                                     bold=run.get('bold', st.get('bold', False)),
                                     italic=run.get('italic', st.get('italic', False)))
            # fill
            fill = st.get('fill')
            if fill:
                set_shape_fill(cell, fill, theme)
            # border (cell 用 XML 设置四边)
            border = st.get('border')
            if border:
                _set_cell_border(cell, border, theme)


# ─────────────────────────── 主流程 ───────────────────────────
class RenderError(RuntimeError):
    """带阶段标记的渲染错误 (借鉴 WASM 的 fail_stage 分层报错设计)。"""

    def __init__(self, stage, message):
        self.stage = stage
        super().__init__(f'[{stage}] {message}')


def render(manifest_path: Path, output_path: Path, debug=False, progress=False):
    if progress:
        print('[progress] 10% 读取 PPTD 项目')
    try:
        pptd_data, pages, root = load_pptd(manifest_path)
    except Exception as e:
        raise RenderError('parse', f'PPTD 解析失败: {e}') from e
    try:
        theme = Theme(pptd_data.get('theme') or {})
    except Exception as e:
        raise RenderError('parse', f'主题解析失败: {e}') from e
    size = pptd_data.get('size') or DEFAULT_SIZE
    page_w, page_h = size[0], size[1]

    prs = Presentation()
    prs.slide_width = Inches(page_w / 72)
    prs.slide_height = Inches(page_h / 72)
    blank = prs.slide_layouts[6]

    for idx, page in enumerate(pages):
        if debug:
            print(f'[page {idx+1}] {page["path"]}')
        if progress:
            print(f'[progress] {15 + int(70 * idx / max(len(pages), 1))}% 渲染第 {idx+1}/{len(pages)} 页')
        try:
            data = page['data']
            slide = prs.slides.add_slide(blank)
            bg = data.get('background')
            if bg:
                bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                                  prs.slide_width, prs.slide_height)
                set_shape_fill(bg_shape, bg, theme)
                bg_shape.line.fill.background()
                bg_shape.shadow.inherit = False
            for el in data.get('elements') or []:
                render_element(slide, el, theme, page_w, page_h, root, debug=debug)
        except RenderError:
            raise
        except Exception as e:
            raise RenderError(f'render-page{idx+1}', f'第 {idx+1} 页渲染失败: {e}') from e

    if progress:
        print('[progress] 85% 保存 PPTX')
    try:
        prs.save(str(output_path))
    except Exception as e:
        raise RenderError('save', f'PPTX 保存失败: {e}') from e

    # ZIP 完整性校验 (借鉴 WASM 的 zip crate 思路: 容器正确性 = PPTX 合法性)
    try:
        import zipfile
        with zipfile.ZipFile(str(output_path)) as zf:
            bad = zf.testzip()
        if bad:
            raise RenderError('verify', f'ZIP 完整性校验失败: 损坏条目 {bad}')
    except RenderError:
        raise
    except Exception as e:
        raise RenderError('verify', f'ZIP 校验异常: {e}') from e

    if progress:
        print('[progress] 100% 完成')
    print(f'saved: {output_path} ({output_path.stat().st_size} bytes)')


def main():
    ap = argparse.ArgumentParser(description='PPTD → PPTX 本地渲染器')
    ap.add_argument('input', help='.pptd 主入口路径')
    ap.add_argument('--output', '-o', required=True, help='输出 .pptx 路径')
    ap.add_argument('--debug', action='store_true', help='打印渲染明细')
    ap.add_argument('--progress', action='store_true', help='打印进度(借鉴 WASM onProgress)')
    args = ap.parse_args()
    try:
        render(Path(args.input), Path(args.output), debug=args.debug, progress=args.progress)
    except RenderError as e:
        print(f'渲染失败 @ {e.stage}: {e}', file=sys.stderr)
        sys.exit(2)


if __name__ == '__main__':
    main()
